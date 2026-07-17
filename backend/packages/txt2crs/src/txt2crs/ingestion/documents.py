# SPDX-License-Identifier: MIT-0

"""DOCX and PPTX text extraction with document/slide boundaries."""

from io import BytesIO

from docx import Document
from pptx import Presentation

from txt2crs.domain.models import InputLocation
from txt2crs.ingestion.errors import EmptyInputError, UnsupportedInputError
from txt2crs.ingestion.models import ExtractedContent, InputPayload


class DocxAdapter:
    """Extract paragraph and table text from a DOCX upload."""

    def extract(self, payload: InputPayload) -> ExtractedContent:
        """Return readable text including table rows."""

        if not isinstance(payload.value, bytes):
            raise UnsupportedInputError("DOCX input must be uploaded as bytes.")
        try:
            document = Document(BytesIO(payload.value))
        except Exception as parsing_error:
            raise UnsupportedInputError(
                "DOCX input is corrupt or unsupported."
            ) from parsing_error

        content_lines = [
            paragraph.text.strip()
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        ]
        for table in document.tables:
            for row in table.rows:
                cell_texts = [cell.text.strip() for cell in row.cells]
                if any(cell_texts):
                    content_lines.append(" | ".join(cell_texts))
        if not content_lines:
            raise EmptyInputError("DOCX contains no extractable text.")
        return ExtractedContent(
            normalized_text="\n".join(content_lines),
            media_type=payload.media_type,
            metadata={"format": "docx"},
            warnings=[],
            locations=[InputLocation(label="Document")],
        )


class PptxAdapter:
    """Extract visible text and preserve 1-based slide locations."""

    def extract(self, payload: InputPayload) -> ExtractedContent:
        """Return slide text in presentation order."""

        if not isinstance(payload.value, bytes):
            raise UnsupportedInputError("PPTX input must be uploaded as bytes.")
        try:
            presentation = Presentation(BytesIO(payload.value))
        except Exception as parsing_error:
            raise UnsupportedInputError(
                "PPTX input is corrupt or unsupported."
            ) from parsing_error

        slide_blocks: list[str] = []
        locations: list[InputLocation] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            shape_texts = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    shape_texts.append(text.strip())
            if not shape_texts:
                continue
            slide_blocks.append("\n".join(shape_texts))
            locations.append(
                InputLocation(label=f"Slide {slide_number}", page=slide_number)
            )
        if not slide_blocks:
            raise EmptyInputError("PPTX contains no extractable text.")
        return ExtractedContent(
            normalized_text="\n\n".join(slide_blocks),
            media_type=payload.media_type,
            metadata={"format": "pptx", "slide_count": len(presentation.slides)},
            warnings=[],
            locations=locations,
        )
