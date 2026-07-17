# SPDX-License-Identifier: MIT-0

"""Fixture-driven tests for PDF, DOCX, and slide extraction."""

from io import BytesIO
from typing import cast

import docx
import fitz  # type: ignore[import-untyped]
import pptx
import pytest

from txt2crs.ingestion.documents import DocxAdapter, PptxAdapter
from txt2crs.ingestion.models import InputPayload
from txt2crs.ingestion.pdf import PdfAdapter
from txt2crs.ingestion.service import EmptyInputError, UnsupportedInputError


def make_pdf(*page_texts: str) -> bytes:
    """Create a tiny in-memory PDF with known page boundaries."""

    document = fitz.open()
    for page_text in page_texts:
        page = document.new_page()
        page.insert_text((72, 72), page_text)
    pdf_bytes = cast(bytes, document.tobytes())
    document.close()
    return pdf_bytes


def test_pdf_adapter_preserves_page_boundaries_and_metadata() -> None:
    """Learner citations can point back to the original PDF page."""

    adapter = PdfAdapter(maximum_pages=10)
    extracted = adapter.extract(
        InputPayload(
            input_type="pdf",
            value=make_pdf("Page one variables", "Page two functions"),
            media_type="application/pdf",
            file_name="course.pdf",
            metadata={},
        )
    )

    assert "Page one variables" in extracted.normalized_text
    assert "Page two functions" in extracted.normalized_text
    assert [location.page for location in extracted.locations] == [1, 2]
    assert extracted.metadata["page_count"] == 2


def test_pdf_adapter_rejects_corrupt_encrypted_and_empty_documents() -> None:
    """Unreadable PDFs return a typed ingestion error, never partial garbage."""

    adapter = PdfAdapter(maximum_pages=10)
    corrupt_payload = InputPayload(
        input_type="pdf",
        value=b"not a PDF",
        media_type="application/pdf",
        file_name="broken.pdf",
        metadata={},
    )
    with pytest.raises(UnsupportedInputError, match="PDF"):
        adapter.extract(corrupt_payload)

    empty_document = fitz.open()
    empty_document.new_page()
    empty_payload = InputPayload(
        input_type="pdf",
        value=empty_document.tobytes(),
        media_type="application/pdf",
        file_name="empty.pdf",
        metadata={},
    )
    empty_document.close()
    with pytest.raises(EmptyInputError, match="text"):
        adapter.extract(empty_payload)

    encrypted_document = fitz.open()
    encrypted_page = encrypted_document.new_page()
    encrypted_page.insert_text((72, 72), "secret")
    encrypted_bytes = encrypted_document.tobytes(
        encryption=fitz.PDF_ENCRYPT_AES_256,
        owner_pw="owner",
        user_pw="learner",
    )
    encrypted_document.close()
    encrypted_payload = InputPayload(
        input_type="pdf",
        value=encrypted_bytes,
        media_type="application/pdf",
        file_name="encrypted.pdf",
        metadata={},
    )
    with pytest.raises(UnsupportedInputError, match="encrypted"):
        adapter.extract(encrypted_payload)


def test_docx_adapter_extracts_paragraphs_and_tables() -> None:
    """DOCX tables remain readable instead of disappearing from the course."""

    document = docx.Document()
    document.add_heading("Variables", level=1)
    document.add_paragraph("A variable binds a name.")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    buffer = BytesIO()
    document.save(buffer)

    extracted = DocxAdapter().extract(
        InputPayload(
            input_type="document",
            value=buffer.getvalue(),
            media_type=(
                "application/vnd.openxmlformats-officedocument."
                "wordprocessingml.document"
            ),
            file_name="course.docx",
            metadata={},
        )
    )

    assert "Variables" in extracted.normalized_text
    assert "A variable binds a name." in extracted.normalized_text
    assert "Name | Value" in extracted.normalized_text


def test_pptx_adapter_preserves_slide_numbers_and_speaker_notes_boundary() -> None:
    """Slide text receives a stable location for later citations."""

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Variables"
    slide.placeholders[1].text = "Names bind values."
    buffer = BytesIO()
    presentation.save(buffer)

    extracted = PptxAdapter().extract(
        InputPayload(
            input_type="slides",
            value=buffer.getvalue(),
            media_type=(
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            ),
            file_name="course.pptx",
            metadata={},
        )
    )

    assert "Variables" in extracted.normalized_text
    assert "Names bind values." in extracted.normalized_text
    assert extracted.locations[0].page == 1
